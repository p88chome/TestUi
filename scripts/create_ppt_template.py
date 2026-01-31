# -*- coding: utf-8 -*-
"""
專業 PPT 模板建立腳本
建立包含多種 Layout 的專業簡報模板
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Deloitte Green Theme
THEME = {
    "primary": RGBColor(0x86, 0xBC, 0x25),      # Deloitte Green
    "secondary": RGBColor(0x00, 0x76, 0xA8),    # Deloitte Blue
    "dark": RGBColor(0x00, 0x00, 0x00),         # Black
    "light": RGBColor(0xFF, 0xFF, 0xFF),        # White
    "gray": RGBColor(0x75, 0x78, 0x7B),         # Gray
    "light_gray": RGBColor(0xF5, 0xF5, 0xF5),   # Light Gray
}


def create_template():
    """建立專業 PPT 模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # ============================================================
    # Layout 0: Title Slide (標題頁)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # 綠色頂部條
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = THEME["primary"]
    top_bar.line.fill.background()
    
    # 主標題區域
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "{{TITLE}}"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = THEME["dark"]
    p.alignment = PP_ALIGN.LEFT
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.8))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "{{SUBTITLE}}"
    p2.font.size = Pt(24)
    p2.font.color.rgb = THEME["gray"]
    p2.alignment = PP_ALIGN.LEFT
    
    # 底部綠色條
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), prs.slide_width, Inches(0.5))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = THEME["primary"]
    bottom_bar.line.fill.background()
    
    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(3), Inches(0.3))
    tf3 = date_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "{{DATE}}"
    p3.font.size = Pt(14)
    p3.font.color.rgb = THEME["light"]
    
    # ============================================================
    # Layout 1: Content Slide (內容頁)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME["primary"]
    header.line.fill.background()
    
    # 頁面標題
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "{{SLIDE_TITLE}}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME["light"]
    
    # 內容區域
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12), Inches(5.3))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "{{CONTENT}}"
    p.font.size = Pt(20)
    p.font.color.rgb = THEME["dark"]
    p.line_spacing = 1.5
    
    # 頁碼
    page_num = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.5), Inches(0.3))
    tf = page_num.text_frame
    p = tf.paragraphs[0]
    p.text = "{{PAGE}}"
    p.font.size = Pt(12)
    p.font.color.rgb = THEME["gray"]
    p.alignment = PP_ALIGN.RIGHT
    
    # ============================================================
    # Layout 2: Section Slide (章節分隔頁)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 全版背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME["primary"]
    bg.line.fill.background()
    
    # 章節標題
    section_title = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(11.7), Inches(1.2))
    tf = section_title.text_frame
    p = tf.paragraphs[0]
    p.text = "{{SECTION_TITLE}}"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = THEME["light"]
    p.alignment = PP_ALIGN.LEFT
    
    # 章節編號
    section_num = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(2), Inches(0.6))
    tf = section_num.text_frame
    p = tf.paragraphs[0]
    p.text = "{{SECTION_NUM}}"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    
    # ============================================================
    # Layout 3: Two Column (雙欄對比頁)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME["primary"]
    header.line.fill.background()
    
    # 頁面標題
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "{{SLIDE_TITLE}}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME["light"]
    
    # 左欄
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.3))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "{{LEFT_CONTENT}}"
    p.font.size = Pt(18)
    p.font.color.rgb = THEME["dark"]
    
    # 右欄
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.7), Inches(5.8), Inches(5.3))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "{{RIGHT_CONTENT}}"
    p.font.size = Pt(18)
    p.font.color.rgb = THEME["dark"]
    
    # 分隔線
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.7), Inches(0.02), Inches(5.3))
    divider.fill.solid()
    divider.fill.fore_color.rgb = THEME["light_gray"]
    divider.line.fill.background()
    
    # ============================================================
    # Layout 4: Bullet Points (項目符號頁)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME["primary"]
    header.line.fill.background()
    
    # 頁面標題
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "{{SLIDE_TITLE}}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME["light"]
    
    # 項目符號內容 - 用多個 bullet point
    bullet_y = 1.8
    for i in range(5):
        # Bullet icon (綠色圓點)
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(bullet_y + i * 1), Inches(0.15), Inches(0.15))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = THEME["primary"]
        bullet.line.fill.background()
        
        # Bullet text
        bullet_text = slide.shapes.add_textbox(Inches(1.0), Inches(bullet_y - 0.1 + i * 1), Inches(11.5), Inches(0.8))
        tf = bullet_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{{{{BULLET_{i+1}}}}}"
        p.font.size = Pt(22)
        p.font.color.rgb = THEME["dark"]
    
    # ============================================================
    # Layout 5: Thank You / End Slide (結束頁)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 綠色頂部條
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = THEME["primary"]
    top_bar.line.fill.background()
    
    # Thank You
    thank_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = THEME["primary"]
    p.alignment = PP_ALIGN.CENTER
    
    # 聯絡資訊
    contact_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "{{CONTACT}}"
    p.font.size = Pt(20)
    p.font.color.rgb = THEME["gray"]
    p.alignment = PP_ALIGN.CENTER
    
    # 底部綠色條
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), prs.slide_width, Inches(0.5))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = THEME["primary"]
    bottom_bar.line.fill.background()
    
    # 儲存模板
    output_path = "app/templates/ppt/deloitte_template.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Template saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_template()
