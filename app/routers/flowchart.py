"""
Flowchart Generator API Router
提供 AI 互動式流程圖產生服務（內控專用版）
- /analyze  : 上傳文件 → AI 產出初始流程圖
- /chat     : 多輪對話修改流程圖
- /export-pptx : 匯出可編輯 PPTX（真正圖形，非圖片）
"""

import io
import json
import re
from datetime import datetime
from typing import List, Optional

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from pydantic import BaseModel

from app.api.deps import get_db, get_current_user
from app.models.user import User
from app.core.config import settings
import httpx

router = APIRouter(prefix="/flowchart", tags=["flowchart"])


# ─── LLM Helper ─────────────────────────────────────────────────────────────

def _get_llm_url_and_headers():
    api_key = settings.AZURE_OPENAI_API_KEY
    endpoint = settings.AZURE_OPENAI_ENDPOINT
    deployment = getattr(settings, "AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4.1")
    api_version = getattr(settings, "AZURE_OPENAI_API_VERSION", "2024-12-01-preview")
    if not api_key:
        raise HTTPException(status_code=500, detail="未設定 Azure OpenAI API Key")
    url = f"{endpoint.rstrip('/')}/openai/deployments/{deployment}/chat/completions?api-version={api_version}"
    headers = {"Content-Type": "application/json", "api-key": api_key}
    return url, headers


# ─── Document Text Extraction ─────────────────────────────────────────────

def _extract_text(filename: str, content: bytes) -> str:
    name_lower = filename.lower()
    if name_lower.endswith((".txt", ".md")):
        return content.decode("utf-8", errors="ignore")
    if name_lower.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"無法解析 DOCX 檔案: {e}")
    if name_lower.endswith(".pdf"):
        try:
            import pdfplumber
            parts = []
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        parts.append(t)
            return "\n".join(parts)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"無法解析 PDF 檔案: {e}")
    return content.decode("utf-8", errors="ignore")


# ─── 【第一步】結構拆解 Prompt ────────────────────────────────────────────────

DECOMPOSE_PROMPT = """你是一位資深的內部控制顧問，請你仔細閱讀以下的作業程序說明（可能是 SOP、稽核文件、Email 或逐字稿），
將其拆解成一個清晰的、結構化的流程步驟清單。

## 輸出格式（只輸出 JSON，不加任何 Markdown 語法或多餘說明）

```json
{
  "title": "流程名稱",
  "summary": "一段話說明這個流程的目的與範圍",
  "steps": [
    {
      "id": "step_1",
      "action": "動作描述，盡量精簡（15字以內）",
      "role": "負責角色或部門（用於泳道分組）",
      "is_decision": false,
      "condition_yes": null,
      "condition_no": null,
      "next_ids": ["step_2"],
      "generates_document": null
    },
    {
      "id": "step_2",
      "action": "主管審核申請單？",
      "role": "部門主管",
      "is_decision": true,
      "condition_yes": "step_3",
      "condition_no": "step_1",
      "next_ids": ["step_3", "step_1"],
      "generates_document": "採購申請單"
    }
  ]
}
```

## 規則
- `is_decision: true` → 這是一個判斷點（是/否分歧），`condition_yes` 和 `condition_no` 必須填入下一步的 step id
- `generates_document` → 如果這個步驟會產生或需要一份特定表單，填入表單名稱，否則填 null
- `next_ids` → 正常情況下填入下一步的 id；如果是決策點，填入所有分支的 id
- 只需要輸出 JSON，不要加任何解釋文字
"""

# ─── 【第二步】圖形渲染 Prompt ────────────────────────────────────────────────

RENDER_PROMPT = """你是一位流程圖繪製專家。根據以下的結構化流程步驟清單，同時產出兩種格式：

1. **nodes/edges JSON**：用於前端 Vue Flow 渲染
2. **Mermaid flowchart 語法**：用於文字視圖顯示

## nodes/edges 格式規則
- 節點 type 只能是：`start` / `end` / `process` / `control` / `document`
- `control` 用於決策點；`document` 用於文件節點（作為側邊分支，不在主流上）
- 每個節點必須有 `lane`（對應步驟的 role）

## Mermaid 格式規則
- 使用 `flowchart TD`（由上至下）
- 決策點用菱形 `{動作}`
- 文件節點用方形 `[文件名]`，用虛線連接到產生它的步驟

## 輸出格式（只輸出 JSON，不加任何 Markdown 語法）

```json
{
  "explanation": "用繁體中文說明這個流程圖的主要結構、涉及角色與關鍵控制點",
  "nodes": [
    {"id": "node_1", "type": "start", "label": "開始", "lane": "申請人"},
    {"id": "node_2", "type": "process", "label": "填寫申請單", "lane": "申請人"},
    {"id": "node_3", "type": "control", "label": "主管審核?", "lane": "部門主管"},
    {"id": "node_4", "type": "document", "label": "採購申請單", "lane": "申請人"},
    {"id": "node_5", "type": "end", "label": "結束", "lane": "採購部"}
  ],
  "edges": [
    {"from": "node_1", "to": "node_2", "label": ""},
    {"from": "node_2", "to": "node_3", "label": ""},
    {"from": "node_2", "to": "node_4", "label": "產出"},
    {"from": "node_3", "to": "node_5", "label": "核准"},
    {"from": "node_3", "to": "node_2", "label": "退回"}
  ],
  "mermaid": "flowchart TD\\n  node_1([開始]):::start --> node_2[填寫申請單]\\n  node_2 --> node_3{主管審核?}\\n  node_2 -.->|產出| node_4[/採購申請單/]\\n  node_3 -->|核准| node_5([結束]):::end\\n  node_3 -->|退回| node_2\\n  classDef start fill:#16a34a,color:#fff\\n  classDef end fill:#dc2626,color:#fff"
}
```
"""


# ─── LLM JSON Parser ─────────────────────────────────────────────────────────

def _parse_llm_json(raw: str) -> dict:
    raw = raw.strip()
    # Strip markdown fences
    if raw.startswith("```"):
        lines = raw.split("\n")
        lines = lines[1:] if lines[0].startswith("```") else lines
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        raw = "\n".join(lines)
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        # Attempt to extract JSON block from prose
        match = re.search(r'\{[\s\S]+\}', raw)
        if match:
            try:
                return json.loads(match.group())
            except:
                pass
        return {
            "explanation": "無法解析 AI 回應的 JSON 格式，請重試。",
            "nodes": [],
            "edges": [],
            "mermaid": ""
        }


# ─── Two-step LLM calls ──────────────────────────────────────────────────────

async def _decompose_to_steps(gateway, user_id: int, doc_text: str, extra: str = "") -> dict:
    """Step 1: Ask LLM to decompose document into structured steps (intermediate layer)."""
    user_content = f"請將以下內容拆解成結構化流程步驟清單：\n\n---\n{doc_text}\n---"
    if extra:
        user_content += f"\n\n補充說明：{extra}"

    messages = [
        {"role": "system", "content": DECOMPOSE_PROMPT},
        {"role": "user", "content": user_content},
    ]
    resp = await gateway.chat(messages=messages, user_id=user_id, temperature=0.1, app_name="Flowchart-Decompose")
    raw = resp["choices"][0]["message"]["content"]
    return _parse_llm_json(raw)


async def _render_steps_to_diagram(gateway, user_id: int, steps_data: dict, modification: str = "") -> dict:
    """Step 2: Convert structured steps into nodes/edges JSON + Mermaid syntax."""
    steps_json = json.dumps(steps_data, ensure_ascii=False, indent=2)
    user_content = f"請根據以下結構化步驟清單，產出流程圖的 nodes/edges JSON 以及 Mermaid 語法：\n\n```json\n{steps_json}\n```"
    if modification:
        user_content += f"\n\n⚠️ 使用者補充修改要求：{modification}\n（請先按修改要求調整步驟結構，再輸出圖形格式）"

    messages = [
        {"role": "system", "content": RENDER_PROMPT},
        {"role": "user", "content": user_content},
    ]
    resp = await gateway.chat(messages=messages, user_id=user_id, temperature=0.1, app_name="Flowchart-Render")
    raw = resp["choices"][0]["message"]["content"]
    return _parse_llm_json(raw)


# ─── Routes ──────────────────────────────────────────────────────────────────

@router.post("/analyze")
async def analyze_document(
    file: UploadFile = File(..., description="內控二階文件 (docx / pdf / txt)"),
    additional_context: Optional[str] = Form(None),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    if not file.filename:
        raise HTTPException(status_code=400, detail="檔案名稱無效")
    content = await file.read()
    if len(content) > 15 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="檔案大小不得超過 15MB")

    doc_text = _extract_text(file.filename, content)
    if not doc_text.strip():
        raise HTTPException(status_code=422, detail="無法從檔案中擷取文字內容")

    if len(doc_text) > 12000:
        doc_text = doc_text[:12000] + "\n\n[...文件內容已截斷...]"

    from app.services.gateway import LLMGateway
    gateway = LLMGateway(db)

    try:
        # ── Step 1: Decompose ──
        steps_data = await _decompose_to_steps(
            gateway, current_user.id, doc_text, additional_context or ""
        )

        # ── Step 2: Render ──
        diagram_data = await _render_steps_to_diagram(gateway, current_user.id, steps_data)

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM API Error: {str(e)}")

    # Build conversation history for chat continuity
    messages = [
        {"role": "system", "content": RENDER_PROMPT},
        {"role": "user", "content": f"[已分析文件: {file.filename}]"},
        {"role": "assistant", "content": json.dumps(diagram_data, ensure_ascii=False)},
    ]

    return {
        "status": "success",
        "explanation": diagram_data.get("explanation", steps_data.get("summary", "")),
        "steps": steps_data,
        "nodes": diagram_data.get("nodes", []),
        "edges": diagram_data.get("edges", []),
        "mermaid": diagram_data.get("mermaid", ""),
        "messages": messages,
        "filename": file.filename,
        "generated_at": datetime.now().isoformat(),
    }


class ChatRequest(BaseModel):
    messages: List[dict]
    user_message: str
    steps: Optional[dict] = None


@router.post("/chat")
async def flowchart_chat(
    request: ChatRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    if not request.user_message.strip():
        raise HTTPException(status_code=400, detail="訊息不可為空")

    from app.services.gateway import LLMGateway
    gateway = LLMGateway(db)

    try:
        if request.steps:
            # ── Two-step path: re-render from updated steps ──
            diagram_data = await _render_steps_to_diagram(
                gateway, current_user.id, request.steps, request.user_message
            )
        else:
            # ── Fallback: single-step legacy path (no steps provided) ──
            messages = list(request.messages)
            messages.append({"role": "user", "content": request.user_message})
            resp = await gateway.chat(messages=messages, user_id=current_user.id, temperature=0.2, app_name="Flowchart-Chat")
            raw = resp["choices"][0]["message"]["content"]
            diagram_data = _parse_llm_json(raw)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM API Error: {str(e)}")

    # Update messages for continuity
    messages = list(request.messages)
    messages.append({"role": "user", "content": request.user_message})
    messages.append({"role": "assistant", "content": json.dumps(diagram_data, ensure_ascii=False)})

    return {
        "status": "success",
        "explanation": diagram_data.get("explanation", ""),
        "steps": request.steps,
        "nodes": diagram_data.get("nodes", []),
        "edges": diagram_data.get("edges", []),
        "mermaid": diagram_data.get("mermaid", ""),
        "messages": messages,
    }


# ─── Text-only Extract (diagram_gen skill) ──────────────────────────────────

class ExtractTextRequest(BaseModel):
    text: str
    title: Optional[str] = None


@router.post("/extract-text")
async def extract_from_text(
    request: ExtractTextRequest,
    current_user: User = Depends(get_current_user),
):
    """純文字輸入 → 走 diagram_gen skill 產初版流程圖
    適合貼單一子作業說明,產第一版給使用者編輯"""
    from app.skills.diagram_gen.extractor import extract_graph
    from app.skills.diagram_gen.converters import to_vueflow_payload
    from app.skills.diagram_gen.mermaid_renderer import to_mermaid

    if not request.text.strip():
        raise HTTPException(status_code=400, detail="文字內容不可為空")

    try:
        graph = extract_graph(request.text)
    except ValueError as e:
        raise HTTPException(status_code=500, detail=f"抽取失敗: {e}")

    payload = to_vueflow_payload(graph)
    mermaid = to_mermaid(graph)

    return {
        "status": "success",
        "explanation": f"已抽取 {len(graph.nodes)} 節點 / {len(graph.edges)} 連線。",
        "title": request.title or graph.title,
        "lanes": payload["lanes"],
        "nodes": payload["nodes"],
        "edges": payload["edges"],
        "mermaid": mermaid,
        "messages": [],
        "generated_at": datetime.now().isoformat(),
    }


# ─── PPTX Export ─────────────────────────────────────────────────────────────

class ExportPptxRequest(BaseModel):
    nodes: List[dict]
    edges: List[dict]
    lanes: Optional[List[dict]] = None  # 可選:前端送的 swimlane 位置
    title: str = "內控流程圖"


@router.post("/export-pptx")
async def export_pptx(
    request: ExportPptxRequest,
    current_user: User = Depends(get_current_user),
):
    """
    將 AI 產出的 nodes/edges 轉換為真正可編輯的 PowerPoint 流程圖。
    若前端提供 nodes[i].x/y 與 lanes[i] 座標(Vue Flow 佈局),
    產出會盡量貼合畫面;否則退回 lane-grouping 簡易版。
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    SLIDE_W = int(prs.slide_width)
    SLIDE_H = int(prs.slide_height)

    # ── 色盤(白底+彩邊, 貼近 Vue Flow 風格) ──
    BORDER = {
        "start":   RGBColor(0x16, 0xa3, 0x4a),
        "end":     RGBColor(0xdc, 0x26, 0x26),
        "process": RGBColor(0x33, 0x41, 0x55),
        "control": RGBColor(0x33, 0x41, 0x55),
    }
    DOC_BORDER = RGBColor(0xea, 0x58, 0x0c)
    EDGE_GRAY = RGBColor(0x47, 0x55, 0x69)
    EDGE_RED = RGBColor(0xdc, 0x26, 0x26)
    TEXT_DARK = RGBColor(0x1e, 0x29, 0x3b)
    LANE_HEADER_BG = RGBColor(0xf1, 0xf5, 0xf9)
    LANE_BORDER = RGBColor(0x94, 0xa3, 0xb8)

    # ── Slide title ──
    tbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.05), SLIDE_W - Inches(0.6), Inches(0.45))
    tf = tbox.text_frame
    tf.text = request.title
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _r = tf.paragraphs[0].runs[0]
    _r.font.size = Pt(16); _r.font.bold = True; _r.font.color.rgb = TEXT_DARK

    # ── 計算座標轉換 (frontend px → slide EMU) ──
    nodes = list(request.nodes)
    lanes = list(request.lanes or [])
    has_layout = all(("x" in n and "y" in n) for n in nodes) and len(nodes) > 0

    TOP_OFFSET = Inches(0.6)
    AREA_W = SLIDE_W - Inches(0.4)
    AREA_H = SLIDE_H - Inches(0.7)

    # Default node box sizes (px, match frontend dagre baseDim)
    DEFAULT_W = {"start": 120, "end": 120, "process": 160, "control": 160}
    DEFAULT_H = {"start": 48, "end": 48, "process": 60, "control": 100}

    def node_px_size(n: dict) -> tuple[int, int]:
        t = n.get("type", "process")
        w = DEFAULT_W.get(t, 160)
        h = DEFAULT_H.get(t, 60)
        return w, h

    if has_layout:
        xs = []
        ys = []
        for n in nodes:
            w, h = node_px_size(n)
            xs.append(n["x"]); xs.append(n["x"] + w + (200 if n.get("docs") else 0))
            ys.append(n["y"]); ys.append(n["y"] + h)
        for l in lanes:
            xs.append(l["x"]); xs.append(l["x"] + l.get("width", 400))
            ys.append(l["y"]); ys.append(l["y"] + l.get("height", 800))
        min_x, max_x = min(xs), max(xs)
        min_y, max_y = min(ys), max(ys)
        scale_x = AREA_W / max(max_x - min_x, 1)
        scale_y = AREA_H / max(max_y - min_y, 1)
        scale = min(scale_x, scale_y) * 0.95

        def to_emu(px_x: float, px_y: float) -> tuple[int, int]:
            return (
                int(Inches(0.2) + (px_x - min_x) * scale),
                int(TOP_OFFSET + (px_y - min_y) * scale),
            )
    else:
        scale = 1.0
        def to_emu(px_x: float, px_y: float) -> tuple[int, int]:
            return (int(Inches(0.2) + px_x), int(TOP_OFFSET + px_y))

    # ── 畫 swimlane 背景 + header (若前端有送) ──
    if has_layout and lanes:
        for lane in lanes:
            lx, ly = to_emu(lane["x"], lane["y"])
            lw = int(lane.get("width", 400) * scale)
            lh = int(lane.get("height", 800) * scale)
            # Lane 外框
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, ly, lw, lh)
            box.fill.background()
            box.line.color.rgb = LANE_BORDER
            box.line.width = Pt(0.75)
            box.shadow.inherit = False
            # Header 條
            hh = int(Inches(0.3))
            hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, ly, lw, hh)
            hdr.fill.solid(); hdr.fill.fore_color.rgb = LANE_HEADER_BG
            hdr.line.color.rgb = LANE_BORDER
            htf = hdr.text_frame
            htf.text = lane.get("label", "")
            hp = htf.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
            hr = hp.runs[0]; hr.font.size = Pt(11); hr.font.bold = True; hr.font.color.rgb = TEXT_DARK

    # ── 畫 nodes ──
    node_boxes: dict[str, tuple[int, int, int, int]] = {}  # id -> (left, top, w, h) of main shape

    for n in nodes:
        t = n.get("type", "process")
        label = n.get("label", "")
        docs = n.get("docs") or []
        px_x = n.get("x", 0); px_y = n.get("y", 0)
        left, top = to_emu(px_x, px_y)
        pw, ph = node_px_size(n)
        w = int(pw * scale); h = int(ph * scale)

        if t == "start" or t == "end":
            shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        elif t == "control":
            shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, w, h)
        else:
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)

        shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
        shp.line.color.rgb = BORDER.get(t, BORDER["process"]); shp.line.width = Pt(1.75)
        tf = shp.text_frame; tf.word_wrap = True; tf.text = label
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = TEXT_DARK

        node_boxes[n["id"]] = (left, top, w, h)

        # ── Doc 貼紙(parallelogram) 放在節點右側 ──
        if docs:
            gap = int(Inches(0.12))
            doc_w = int(160 * scale)
            doc_h = int((46 + 22 * len(docs)) * scale)
            doc_left = left + w + gap
            doc_top = top + h // 2 - doc_h // 2
            dshp = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, doc_left, doc_top, doc_w, doc_h)
            dshp.fill.solid(); dshp.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
            dshp.line.color.rgb = DOC_BORDER; dshp.line.width = Pt(1.5)
            dtf = dshp.text_frame; dtf.word_wrap = True
            dtf.text = docs[0]
            dtf.paragraphs[0].alignment = PP_ALIGN.CENTER
            dr = dtf.paragraphs[0].runs[0]
            dr.font.size = Pt(8); dr.font.color.rgb = TEXT_DARK
            for extra in docs[1:]:
                p2 = dtf.add_paragraph()
                p2.text = extra; p2.alignment = PP_ALIGN.CENTER
                er = p2.runs[0]; er.font.size = Pt(8); er.font.color.rgb = TEXT_DARK

    # ── 畫 edges (elbow connector + N 紅虛線) ──
    for e in request.edges:
        fb = node_boxes.get(e.get("from", ""))
        tb = node_boxes.get(e.get("to", ""))
        if not fb or not tb:
            continue
        fl, ft, fw, fh = fb
        tl, tt, tw, th = tb
        raw_label = str(e.get("label", "") or "").strip()
        is_no = raw_label.upper() in ("N", "NO") or raw_label == "否"
        is_yes = raw_label.upper() in ("Y", "YES") or raw_label == "是"

        # 起終錨點: N 走左側, 其他底/頂
        if is_no:
            x1 = fl
            y1 = ft + fh // 2
            x2 = tl
            y2 = tt + th // 2
        else:
            x1 = fl + fw // 2
            y1 = ft + fh
            x2 = tl + tw // 2
            y2 = tt

        # 2 = ELBOW_CONNECTOR
        conn = slide.shapes.add_connector(2, x1, y1, x2, y2)
        if is_no:
            conn.line.color.rgb = EDGE_RED
            conn.line.width = Pt(1.25)
        else:
            conn.line.color.rgb = EDGE_GRAY
            conn.line.width = Pt(1.5)

        # 加箭頭 + N 虛線樣式 (OXML manipulation)
        from pptx.oxml.ns import qn
        ln = conn.line._get_or_add_ln()
        # Dashed for N
        if is_no and ln.find(qn("a:prstDash")) is None:
            dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
            ln.append(dash)
        # 箭頭 (只加一次)
        if ln.find(qn("a:tailEnd")) is None:
            tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "h": "sm"})
            ln.append(tail)

        # label (主幹 Y 不顯示)
        if raw_label and not is_yes:
            mid_x = (x1 + x2) // 2
            mid_y = (y1 + y2) // 2
            lbl = slide.shapes.add_textbox(
                mid_x - Inches(0.35), mid_y - Inches(0.12), Inches(0.7), Inches(0.25)
            )
            ltf = lbl.text_frame; ltf.text = raw_label
            lp = ltf.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
            lr = lp.runs[0]; lr.font.size = Pt(8); lr.font.bold = True
            lr.font.color.rgb = EDGE_RED if is_no else TEXT_DARK

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    filename = f"flowchart_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )
