# -*- coding: utf-8 -*-
"""
PPTX Generator - 專業模板版
使用預製的專業模板生成高品質簡報
"""

import os
from typing import List, Dict, Any, Optional
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# 模板檔案路徑
TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "templates", "ppt")


# 色彩主題
THEMES = {
    "deloitte": {
        "primary": RGBColor(0x86, 0xBC, 0x25),
        "secondary": RGBColor(0x00, 0x76, 0xA8),
        "dark": RGBColor(0x00, 0x00, 0x00),
        "light": RGBColor(0xFF, 0xFF, 0xFF),
        "gray": RGBColor(0x75, 0x78, 0x7B),
    },
    "professional": {
        "primary": RGBColor(0x1E, 0x3A, 0x5F),
        "secondary": RGBColor(0x3D, 0x5A, 0x80),
        "dark": RGBColor(0x0D, 0x1B, 0x2A),
        "light": RGBColor(0xFF, 0xFF, 0xFF),
        "gray": RGBColor(0x6B, 0x70, 0x7A),
    },
    "creative": {
        "primary": RGBColor(0x8B, 0x5C, 0xF6),
        "secondary": RGBColor(0xF9, 0x73, 0x16),
        "dark": RGBColor(0x1F, 0x2A, 0x37),
        "light": RGBColor(0xFF, 0xFF, 0xFF),
        "gray": RGBColor(0x64, 0x74, 0x8B),
    },
    "minimal": {
        "primary": RGBColor(0x37, 0x37, 0x37),
        "secondary": RGBColor(0x6B, 0x6B, 0x6B),
        "dark": RGBColor(0x00, 0x00, 0x00),
        "light": RGBColor(0xFF, 0xFF, 0xFF),
        "gray": RGBColor(0x9E, 0x9E, 0x9E),
    }
}


def execute(input_data: dict) -> dict:
    """
    執行 PPTX 生成
    """
    if not PPTX_AVAILABLE:
        return {
            "status": "error",
            "message": "python-pptx 未安裝。請執行: pip install python-pptx"
        }
    
    operation = input_data.get("operation", "create_from_outline")
    
    try:
        if operation == "create_from_outline":
            return _create_from_outline(input_data)
        elif operation == "add_slide":
            return _add_slide(input_data)
        else:
            return {"status": "error", "message": f"不支援的操作: {operation}"}
    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"status": "error", "message": str(e)}


def _create_from_outline(input_data: dict) -> dict:
    """從大綱建立簡報 - 使用專業模板風格"""
    title = input_data.get("title", "簡報標題")
    subtitle = input_data.get("subtitle", "")
    slides = input_data.get("slides", [])
    theme_name = input_data.get("theme", "deloitte")
    output_path = input_data.get("output_path")
    
    theme = THEMES.get(theme_name, THEMES["deloitte"])
    
    # 建立新簡報
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 添加標題頁
    _add_title_slide(prs, title, subtitle, theme)
    
    # 添加內容投影片
    for i, slide_data in enumerate(slides):
        slide_layout = slide_data.get("layout", "bullet")
        slide_title = slide_data.get("title", f"投影片 {i+1}")
        content = slide_data.get("content", [])
        
        if slide_layout == "section":
            _add_section_slide(prs, slide_title, i+1, theme)
        elif slide_layout == "two_column":
            left = content[:len(content)//2] if content else []
            right = content[len(content)//2:] if content else []
            _add_two_column_slide(prs, slide_title, left, right, theme)
        else:  # bullet or content
            _add_bullet_slide(prs, slide_title, content, theme)
    
    # 添加結束頁
    _add_thank_you_slide(prs, theme)
    
    # 儲存
    if not output_path:
        output_path = os.path.join(os.getcwd(), f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
    
    prs.save(output_path)
    
    return {
        "status": "success",
        "file_path": os.path.abspath(output_path),
        "slide_count": len(prs.slides),
        "theme": theme_name
    }


def _add_title_slide(prs, title: str, subtitle: str, theme: dict):
    """添加專業標題頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # 頂部綠色條
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = theme["primary"]
    top_bar.line.fill.background()
    
    # 主標題
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = theme["dark"]
    p.alignment = PP_ALIGN.LEFT
    
    # 副標題
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.8))
        tf2 = subtitle_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = theme["gray"]
        p2.alignment = PP_ALIGN.LEFT
    
    # 底部條
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), prs.slide_width, Inches(0.5))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = theme["primary"]
    bottom_bar.line.fill.background()
    
    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(3), Inches(0.3))
    tf3 = date_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = datetime.now().strftime("%Y年%m月%d日")
    p3.font.size = Pt(14)
    p3.font.color.rgb = theme["light"]


def _add_bullet_slide(prs, title: str, content: List[str], theme: dict):
    """添加項目符號頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = theme["primary"]
    header.line.fill.background()
    
    # 頁面標題
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["light"]
    
    # 項目符號
    bullet_y = 1.8
    line_height = 1.0
    
    for i, item in enumerate(content[:5]):  # 最多 5 個項目
        # 綠色圓點
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            Inches(0.6), 
            Inches(bullet_y + i * line_height + 0.1), 
            Inches(0.15), 
            Inches(0.15)
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = theme["primary"]
        bullet.line.fill.background()
        
        # 文字
        text_box = slide.shapes.add_textbox(
            Inches(1.0), 
            Inches(bullet_y + i * line_height), 
            Inches(11.5), 
            Inches(0.8)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item if isinstance(item, str) else str(item)
        p.font.size = Pt(22)
        p.font.color.rgb = theme["dark"]
    
    # 頁碼
    _add_page_number(slide, len(prs.slides), theme)


def _add_section_slide(prs, title: str, section_num: int, theme: dict):
    """添加章節分隔頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 全版背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["primary"]
    bg.line.fill.background()
    
    # 章節編號
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(2), Inches(0.6))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_num}" if section_num < 10 else str(section_num)
    p.font.size = Pt(24)
    p.font.color.rgb = theme["light"]
    
    # 章節標題
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = theme["light"]
    p.alignment = PP_ALIGN.LEFT


def _add_two_column_slide(prs, title: str, left_content: List[str], right_content: List[str], theme: dict):
    """添加雙欄對比頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = theme["primary"]
    header.line.fill.background()
    
    # 頁面標題
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["light"]
    
    # 左欄
    left_text = "\n• ".join(left_content) if left_content else ""
    if left_text:
        left_text = "• " + left_text
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.3))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_text
    p.font.size = Pt(18)
    p.font.color.rgb = theme["dark"]
    p.line_spacing = 1.5
    
    # 分隔線
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.7), Inches(0.02), Inches(5.3))
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    divider.line.fill.background()
    
    # 右欄
    right_text = "\n• ".join(right_content) if right_content else ""
    if right_text:
        right_text = "• " + right_text
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.7), Inches(5.8), Inches(5.3))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_text
    p.font.size = Pt(18)
    p.font.color.rgb = theme["dark"]
    p.line_spacing = 1.5
    
    # 頁碼
    _add_page_number(slide, len(prs.slides), theme)


def _add_thank_you_slide(prs, theme: dict):
    """添加結束頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 頂部條
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = theme["primary"]
    top_bar.line.fill.background()
    
    # Thank You
    thank_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = theme["primary"]
    p.alignment = PP_ALIGN.CENTER
    
    # 聯絡資訊
    contact_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(20)
    p.font.color.rgb = theme["gray"]
    p.alignment = PP_ALIGN.CENTER
    
    # 底部條
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), prs.slide_width, Inches(0.5))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = theme["primary"]
    bottom_bar.line.fill.background()


def _add_page_number(slide, page_num: int, theme: dict):
    """添加頁碼"""
    prs = slide.part.package.presentation_part.presentation
    page_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.5), Inches(0.3))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(12)
    p.font.color.rgb = theme["gray"]
    p.alignment = PP_ALIGN.RIGHT


def _add_slide(input_data: dict) -> dict:
    """為現有簡報添加投影片"""
    file_path = input_data.get("file_path")
    slide_data = input_data.get("slide")
    theme_name = input_data.get("theme", "deloitte")
    
    if not file_path or not os.path.exists(file_path):
        return {"status": "error", "message": "找不到簡報檔案"}
    
    theme = THEMES.get(theme_name, THEMES["deloitte"])
    prs = Presentation(file_path)
    
    layout = slide_data.get("layout", "bullet")
    title = slide_data.get("title", "新投影片")
    content = slide_data.get("content", [])
    
    if layout == "section":
        _add_section_slide(prs, title, len(prs.slides), theme)
    elif layout == "two_column":
        left = content[:len(content)//2] if content else []
        right = content[len(content)//2:] if content else []
        _add_two_column_slide(prs, title, left, right, theme)
    else:
        _add_bullet_slide(prs, title, content, theme)
    
    prs.save(file_path)
    
    return {
        "status": "success",
        "file_path": file_path,
        "slide_count": len(prs.slides)
    }
